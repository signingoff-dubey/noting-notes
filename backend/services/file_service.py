"""
File text-extraction service.
Supports PDF, DOCX, XLSX, PPTX. All blocking I/O is offloaded via asyncio.to_thread.
"""

import asyncio
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


# ──────────────── Sync extractors (run in thread pool) ────────────────

def _extract_pdf(file_path: str) -> str:
    """Extract text from a PDF using PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.error("PyMuPDF not installed — cannot extract PDF text")
        return ""

    text_parts: list[str] = []
    try:
        doc = fitz.open(file_path)
        for page in doc:
            page_text = page.get_text()
            if page_text:
                text_parts.append(page_text)
        doc.close()
    except Exception as exc:
        logger.warning("PDF extraction failed for %s: %s", file_path, exc)
    return "\n".join(text_parts)


def _extract_docx(file_path: str) -> str:
    """Extract text from a DOCX file using python-docx (paragraphs + tables)."""
    try:
        from docx import Document
    except ImportError:
        logger.error("python-docx not installed — cannot extract DOCX text")
        return ""

    text_parts: list[str] = []
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_cells:
                    text_parts.append("\t".join(row_cells))
    except Exception as exc:
        logger.warning("DOCX extraction failed for %s: %s", file_path, exc)
    return "\n".join(text_parts)


def _extract_xlsx(file_path: str) -> str:
    """Extract text from an XLSX file using openpyxl (all sheets, all cells)."""
    try:
        import openpyxl
    except ImportError:
        logger.error("openpyxl not installed — cannot extract XLSX text")
        return ""

    text_parts: list[str] = []
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    text_parts.append("\t".join(cells))
        wb.close()
    except Exception as exc:
        logger.warning("XLSX extraction failed for %s: %s", file_path, exc)
    return "\n".join(text_parts)


def _extract_pptx(file_path: str) -> str:
    """Extract text from a PPTX file using python-pptx (all slide shapes)."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed — cannot extract PPTX text")
        return ""

    text_parts: list[str] = []
    try:
        prs = Presentation(file_path)
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                text_parts.append(f"[Slide {slide_idx}]")
                text_parts.extend(slide_texts)
    except Exception as exc:
        logger.warning("PPTX extraction failed for %s: %s", file_path, exc)
    return "\n".join(text_parts)


# ──────────────── Extension → MIME fallback map ────────────────

_EXT_TO_MIME: dict[str, str] = {
    ".pdf":  "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc":  "application/msword",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".xls":  "application/vnd.ms-excel",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppt":  "application/vnd.ms-powerpoint",
}


def _resolve_mime(file_path: str, mime_type: str) -> str:
    """If mime_type is generic/empty, infer from file extension."""
    if mime_type and mime_type not in ("application/octet-stream", ""):
        return mime_type
    ext = Path(file_path).suffix.lower()
    return _EXT_TO_MIME.get(ext, mime_type)


# ──────────────── Public async API ────────────────

async def extract_text(file_path: str, mime_type: str) -> str:
    """
    Extract plain text from *file_path*.

    Supported types:
      - PDF  → PyMuPDF
      - DOCX → python-docx
      - XLSX → openpyxl
      - PPTX → python-pptx
      - Other → returns ""

    All blocking I/O is delegated to asyncio.to_thread.
    """
    resolved_mime = _resolve_mime(file_path, mime_type)

    if resolved_mime == "application/pdf":
        return await asyncio.to_thread(_extract_pdf, file_path)

    if resolved_mime in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        return await asyncio.to_thread(_extract_docx, file_path)

    if resolved_mime in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        return await asyncio.to_thread(_extract_xlsx, file_path)

    if resolved_mime in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        return await asyncio.to_thread(_extract_pptx, file_path)

    # Images or unsupported types — return empty string
    logger.debug("extract_text: unsupported mime_type=%s for %s", resolved_mime, file_path)
    return ""
